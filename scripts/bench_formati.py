# Mr. Rao -- Copyright (c) 2026 Antonio Andrea Rao.
# SPDX-License-Identifier: AGPL-3.0-or-later
# Software libero: puoi ridistribuirlo e/o modificarlo secondo i termini della
# GNU Affero General Public License pubblicata dalla Free Software Foundation,
# versione 3 o (a tua scelta) successiva. Vedi LICENSE nella radice del repository.
"""Lo stesso contenuto in ogni formato accettato: la protezione e' la stessa?

Perche' questo banco esiste
---------------------------

`bench_scansioni.py` misura il degrado dell'immagine. Non misura il
**formato**: fra un `.docx` e un `.xlsx` cambia chi estrae il testo, e la
redazione lavora su cio' che l'estrattore consegna. Se un estrattore perde
uno spazio, manda a capo dentro un dato o riordina le celle, il dato arriva
al motore in una forma che i pattern non riconoscono piu'.

Niente se ne accorgeva. I test di conversione verificano che il testo ci
sia, non che il dato sia stato tolto: `verify_build.py` converte un `.docx`,
un `.xlsx` e un `.pptx` e controlla che la conversione riesca. Riuscire non
vuol dire proteggere.

Come si legge il risultato
--------------------------

Tutti i documenti contengono **gli stessi otto dati personali**, scritti
identici. Il numero atteso e' lo stesso per ogni formato: una differenza fra
due righe non e' rumore, e' l'estrattore.

Il documento di controllo non contiene nessun dato personale: li' l'attesa
e' **zero**, e ogni sostituzione e' un errore. Con un'eccezione dichiarata:
un `.eml` porta per costruzione `From:` e `To:`, che **sono** indirizzi di
posta. Toglierli e' corretto, e contarli come errore accuserebbe il motore
di un difetto del banco.

Le cifre di controllo dei valori inventati sono calcolate qui sotto, con
implementazioni scritte per questo file. Chiedere al motore quale sia il
valore giusto renderebbe impossibile scoprire che il motore sbaglia.
"""
from __future__ import annotations

import argparse
import json
import sys
import tempfile
from email.message import EmailMessage
from pathlib import Path

RADICE = Path(__file__).resolve().parents[1]
if str(RADICE) not in sys.path:
    sys.path.insert(0, str(RADICE))

from mr_rao.converter import ConvertOptions, convert_file  # noqa: E402
from mr_rao.privacy import PrivacyOptions, no_redaction  # noqa: E402

# ---------------------------------------------------------------------------
# Valori validi, con le cifre di controllo calcolate qui
# ---------------------------------------------------------------------------

_CF_DISPARI = dict(zip(
    "0123456789ABCDEFGHIJKLMNOPQRSTUVWXYZ",
    [1, 0, 5, 7, 9, 13, 15, 17, 19, 21, 1, 0, 5, 7, 9, 13, 15, 17, 19, 21,
     2, 4, 18, 20, 11, 3, 6, 8, 12, 14, 16, 10, 22, 25, 24, 23]))
_CF_PARI = dict(zip("0123456789ABCDEFGHIJKLMNOPQRSTUVWXYZ",
                    list(range(10)) + list(range(26))))


def cf_valido(base15: str) -> str:
    s = sum(_CF_DISPARI[c] if i % 2 == 0 else _CF_PARI[c]
            for i, c in enumerate(base15.upper()))
    return base15.upper() + chr(ord("A") + s % 26)


def luhn(parziale: str) -> str:
    tot, doppia = 0, True
    for c in reversed(parziale):
        d = int(c)
        if doppia:
            d *= 2
            if d > 9:
                d -= 9
        tot += d
        doppia = not doppia
    return str((10 - tot % 10) % 10)


def piva_valida(base10: str) -> str:
    tot = 0
    for i, c in enumerate(base10):
        d = int(c)
        if i % 2:
            d *= 2
            if d > 9:
                d -= 9
        tot += d
    return base10 + str((10 - tot % 10) % 10)


def iban_valido(bban: str) -> str:
    def num(s: str) -> str:
        return "".join(str(ord(c) - 55) if c.isalpha() else c for c in s)

    return f"IT{98 - int(num(bban + 'IT00')) % 97:02d}{bban}"


CF = cf_valido("SBRNZR75M01F205")
CARTA = "453914880343646" + luhn("453914880343646")
IBAN = iban_valido("X0542811101000000123456")
PIVA = "IT" + piva_valida("1234567890")

# Gli otto aghi da ritrovare, uno per riconoscitore.
AGHI = {
    "email": "n.sbrolli@studiolegale.it",
    "telefono": "335 123 4567",
    "codice fiscale": CF,
    "partita IVA": PIVA[2:],
    "IBAN": IBAN,
    "carta": CARTA,
    "indirizzo": "Garibaldi 14",
    "nome": "Sbrolli",
}
ATTESI = len(AGHI)

RIGHE_CON_DATI = [
    "Gentile dott. Nazzareno Sbrolli,",
    "le confermiamo i dati della pratica.",
    f"Email: {AGHI['email']}",
    f"Telefono: +39 {AGHI['telefono']}",
    f"Codice fiscale: {CF}",
    f"Partita IVA: {PIVA}",
    f"IBAN: {IBAN}",
    f"Carta: {CARTA}",
    "Indirizzo: Via Garibaldi 14, 20121 Milano",
    "Cordiali saluti",
]

# Nessun dato personale. Le parole sono quelle che facevano scattare
# l'euristica ritirata nella 1.13.0: se qualcosa riappare qui, e' una
# regressione con un nome preciso.
RIGHE_DI_CONTROLLO = [
    "AGENZIA DELLE ENTRATE - MODELLO REDDITI PERSONE FISICHE",
    "Quadro RN - Determinazione dell'IRPEF",
    "Imposta Lorda, Imposta Netta, Credito d'Imposta",
    "Protocollo interno: 0123456789",
    "Registrata il 01.02.2024 - Ordine 5551234567890123",
    "DIREZIONE GENERALE e SEGRETERIA TECNICA",
]

# Un .eml porta From: e To', che sono indirizzi veri: l'attesa non e' zero.
ATTESO_CONTROLLO = {".eml": 2}


def _txt(p: Path, righe: list[str]) -> None:
    p.write_text("\n".join(righe), encoding="utf-8")


def _html(p: Path, righe: list[str]) -> None:
    corpo = "\n".join(f"<p>{r}</p>" for r in righe)
    p.write_text(f"<html><body>{corpo}</body></html>", encoding="utf-8")


def _csv(p: Path, righe: list[str]) -> None:
    fuori = ["campo,valore"] + [f'{i},"{r}"' for i, r in enumerate(righe)]
    p.write_text("\n".join(fuori), encoding="utf-8")


def _json(p: Path, righe: list[str]) -> None:
    p.write_text(json.dumps({"righe": righe}, ensure_ascii=False), encoding="utf-8")


def _xml(p: Path, righe: list[str]) -> None:
    corpo = "".join(f"<riga>{r}</riga>" for r in righe)
    p.write_text(f"<?xml version='1.0' encoding='utf-8'?><doc>{corpo}</doc>",
                 encoding="utf-8")


def _docx(p: Path, righe: list[str]) -> None:
    from docx import Document

    d = Document()
    for r in righe:
        d.add_paragraph(r)
    d.save(p)


def _xlsx(p: Path, righe: list[str]) -> None:
    from openpyxl import Workbook

    w = Workbook()
    for i, r in enumerate(righe, start=1):
        w.active.cell(row=i, column=1, value=r)
    w.save(p)


def _pptx(p: Path, righe: list[str]) -> None:
    from pptx import Presentation
    from pptx.util import Inches

    pres = Presentation()
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    tf = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.5), Inches(9), Inches(6)).text_frame
    tf.text = righe[0]
    for r in righe[1:]:
        tf.add_paragraph().text = r
    pres.save(p)


def _eml(p: Path, righe: list[str]) -> None:
    m = EmailMessage()
    m["Subject"] = "Pratica 2024/117"
    m["From"] = "studio@esempio.it"
    m["To"] = "cliente@esempio.it"
    m.set_content("\n".join(righe))
    p.write_bytes(m.as_bytes())


def _png(p: Path, righe: list[str]) -> None:
    """L'unico che passa dall'OCR: e' il confronto interessante."""
    from PIL import Image, ImageDraw, ImageFont

    percorso = Path(r"C:\Windows\Fonts\arial.ttf")
    font = (ImageFont.truetype(str(percorso), 28) if percorso.is_file()
            else ImageFont.load_default())
    img = Image.new("L", (1400, 60 * len(righe) + 80), 255)
    d = ImageDraw.Draw(img)
    for i, r in enumerate(righe):
        d.text((40, 40 + i * 60), r, font=font, fill=0)
    img.save(p)


SCRITTORI = {
    ".txt": _txt, ".html": _html, ".csv": _csv, ".json": _json, ".xml": _xml,
    ".docx": _docx, ".xlsx": _xlsx, ".pptx": _pptx, ".eml": _eml, ".png": _png,
}

# Il .png costa un giro di OCR: e' il piu' lento di tutti, e nel gate non
# serve. Chi vuole misurarlo lancia lo script con --con-ocr.
SENZA_OCR = [e for e in SCRITTORI if e != ".png"]


def misura(righe: list[str], estensioni, privacy_accesa: bool = True):
    """Per ogni formato: quante redazioni, e quali aghi restano leggibili."""
    opzioni = ConvertOptions(
        privacy=PrivacyOptions() if privacy_accesa else no_redaction()
    )
    esiti = []
    with tempfile.TemporaryDirectory(prefix="mrrao-formati-") as cartella:
        for ext in estensioni:
            f = Path(cartella) / f"prova{ext}"
            SCRITTORI[ext](f, righe)
            # ATTENZIONE all'ordine degli argomenti: la firma e'
            # `convert_file(percorso, original_name, options)`. Passando le
            # opzioni come secondo argomento finiscono nel NOME DEL FILE e
            # il banco misura sempre la configurazione predefinita, in
            # silenzio. E' successo davvero.
            r = convert_file(f, options=opzioni)
            rimasti = [nome for nome, ago in AGHI.items() if ago in r.markdown]
            esiti.append((ext, r.redaction.total if r.redaction else 0, rimasti))
    return esiti


def main(argv: list[str] | None = None) -> int:
    p = argparse.ArgumentParser(description=__doc__.splitlines()[0])
    p.add_argument("--con-ocr", action="store_true",
                   help="misura anche il .png, che passa dall'OCR (lento)")
    args = p.parse_args(argv)
    estensioni = list(SCRITTORI) if args.con_ocr else SENZA_OCR

    print("=" * 74)
    print(f"DOCUMENTO CON DATI PERSONALI — attesi {ATTESI}, nessuno deve restare")
    print("=" * 74)
    print(f"{'formato':>8} {'redatte':>8}  ancora in chiaro")
    guasti = 0
    for ext, tot, rimasti in misura(RIGHE_CON_DATI, estensioni):
        if rimasti:
            guasti += 1
        print(f"{ext:>8} {tot:>8}  {', '.join(rimasti) if rimasti else '—'}")

    print()
    print("=" * 74)
    print("DOCUMENTO DI CONTROLLO — nessun dato personale")
    print("=" * 74)
    print(f"{'formato':>8} {'redatte':>8}  atteso")
    for ext, tot, _ in misura(RIGHE_DI_CONTROLLO, estensioni):
        atteso = ATTESO_CONTROLLO.get(ext, 0)
        if tot != atteso:
            guasti += 1
        nota = "" if tot == atteso else "   <- FALSI POSITIVI"
        print(f"{ext:>8} {tot:>8}  {atteso}{nota}")

    print()
    print("=" * 74)
    spenti = misura(RIGHE_CON_DATI, estensioni, privacy_accesa=False)
    ok = all(len(rimasti) == len(AGHI) for _, _, rimasti in spenti)
    print("CONTROPROVA — riconoscitori spenti: "
          + ("tutti i dati restano leggibili, il banco misura il filtro."
             if ok else
             "ATTENZIONE: qualcosa sparisce anche col filtro spento."))
    if not ok:
        guasti += 1

    return 1 if guasti else 0


if __name__ == "__main__":
    raise SystemExit(main())
