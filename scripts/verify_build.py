# Mr. Rao -- Copyright (c) 2026 Antonio Andrea Rao.
# SPDX-License-Identifier: AGPL-3.0-or-later
# Software libero: puoi ridistribuirlo e/o modificarlo secondo i termini della
# GNU Affero General Public License pubblicata dalla Free Software Foundation,
# versione 3 o (a tua scelta) successiva. Vedi LICENSE nella radice del repository.
"""Avvia l'eseguibile appena costruito e verifica che serva davvero.

Un build che finisce con codice di uscita zero non dice niente su cosa
succede quando si fa doppio clic. E' gia' successo di produrre 390 MB che
aprivano una finestra nera e si chiudevano: PyInstaller aveva incluso un
runtime hook che moriva su un file mancante, prima ancora di arrivare al
nostro codice. Se ne accorse un essere umano avviandolo, non il build.

Da qui in poi se ne accorge il build. I controlli, nell'ordine in cui
fallirebbero:

1. il processo resta vivo dopo l'avvio;
2. /api/health risponde, con la versione giusta e frozen=True;
3. una conversione vera produce testo — e non di un .txt, ma di **tutti e
   tre** i formati Office che dipendono da una libreria opzionale;
4. l'anonimizzazione lavora davvero sul testo estratto;
5. l'icona spedita e' quella del repository.

Perche' tutti e tre e non solo il .docx, che c'era gia'. Ogni formato
Office ha una libreria diversa dietro — mammoth per Word, pandas piu'
openpyxl per Excel, python-pptx per PowerPoint — e per tre versioni quelle
librerie sono finite nel pacchetto **per caso**: erano nel venv di sviluppo
da un'installazione precedente, non nell'elenco delle dipendenze. Ha
funzionato tutto, e nessuno lo stava controllando.

Con un formato solo sotto esame, il giorno che una di quelle librerie non
viene piu' impacchettata il build accetta il pacchetto e se ne accorge chi
lo usa: «nessun testo riconoscibile», con la colpa apparente al documento.

Uso:  python scripts/verify_build.py dist/MrRao-Portable/app/MrRao.exe
"""
from __future__ import annotations

import io
import json
import os
import re
import socket
import subprocess
import sys
import time
import urllib.error
import urllib.request
import uuid
import zipfile
from pathlib import Path

AVVIO_MAX_S = 90
PASSO_S = 1.0


def porta_libera() -> int:
    with socket.socket() as s:
        s.bind(("127.0.0.1", 0))
        return s.getsockname()[1]


def segnaposto_presente(md: str, etichetta: str) -> bool:
    """`{{EMAIL}}` oppure `{{EMAIL_1}}`: la numerazione non e' un'assenza.

    Il controllo cercava la forma piatta alla lettera. Da quando i
    segnaposto escono numerati, `{{EMAIL_1}}` nel testo faceva dire al
    build che l'anonimizzazione mancava: il pacchetto era buono, il metro
    era vecchio. Cambiare il prodotto per far tacere il metro sarebbe
    stato il modo sbagliato di leggerlo.

    Restano due modi di fallire, e sono quelli che contano: il segnaposto
    non c'e' affatto (il dato e' uscito in chiaro), oppure c'e' ma in una
    forma che non riconosciamo -- per esempio col marcatore interno
    ancora attaccato, che non deve mai arrivare a chi legge.
    """
    return re.search(r"\{\{" + etichetta + r"(?:_\d+)?\}\}", md) is not None


def docx_di_prova(testo: str) -> bytes:
    doc = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<w:document xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">'
        f"<w:body><w:p><w:r><w:t>{testo}</w:t></w:r></w:p></w:body></w:document>"
    )
    rels = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
        '<Relationship Id="rId1" Target="word/document.xml" '
        'Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument"/>'
        "</Relationships>"
    )
    types = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">'
        '<Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>'
        '<Default Extension="xml" ContentType="application/xml"/>'
        '<Override PartName="/word/document.xml" ContentType="application/vnd.openxmlformats-'
        'officedocument.wordprocessingml.document.main+xml"/></Types>'
    )
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w") as z:
        z.writestr("[Content_Types].xml", types)
        z.writestr("_rels/.rels", rels)
        z.writestr("word/document.xml", doc)
    return buf.getvalue()


def xlsx_di_prova(testo: str) -> bytes:
    """Un .xlsx vero, scritto da openpyxl.

    Costruirlo a mano come il .docx non conviene: openpyxl vuole un
    pacchetto piu' completo, e un file appena fuori specifica farebbe
    fallire la verifica per il motivo sbagliato.

    Se openpyxl manca nel venv di build l'errore e' quello giusto: senza,
    il pacchetto uscirebbe comunque, con Excel rotto dentro.
    """
    from openpyxl import Workbook

    cartella = Workbook()
    foglio = cartella.active
    foglio["A1"] = "Riferimento"
    foglio["B1"] = testo
    buf = io.BytesIO()
    cartella.save(buf)
    return buf.getvalue()


def pptx_di_prova(testo: str) -> bytes:
    """Un .pptx vero, scritto da python-pptx.

    A mano sarebbe anche peggio del .xlsx: servono slide master e layout.
    Vale lo stesso ragionamento sull'import mancante.
    """
    from pptx import Presentation

    presentazione = Presentation()
    diapositiva = presentazione.slides.add_slide(presentazione.slide_layouts[5])
    diapositiva.shapes.title.text = testo
    buf = io.BytesIO()
    presentazione.save(buf)
    return buf.getvalue()


def post_multipart(url: str, campi: dict[str, str], nome: str, contenuto: bytes) -> dict:
    confine = "----MrRaoVerify" + uuid.uuid4().hex
    corpo = bytearray()
    for k, v in campi.items():
        corpo += f"--{confine}\r\n".encode()
        corpo += f'Content-Disposition: form-data; name="{k}"\r\n\r\n{v}\r\n'.encode()
    corpo += f"--{confine}\r\n".encode()
    corpo += (
        f'Content-Disposition: form-data; name="file"; filename="{nome}"\r\n'
        f"Content-Type: application/octet-stream\r\n\r\n"
    ).encode()
    corpo += contenuto + b"\r\n"
    corpo += f"--{confine}--\r\n".encode()

    req = urllib.request.Request(
        url,
        data=bytes(corpo),
        headers={
            "Content-Type": f"multipart/form-data; boundary={confine}",
            "Origin": url.rsplit("/api/", 1)[0],
        },
    )
    with urllib.request.urlopen(req, timeout=120) as r:
        return json.loads(r.read().decode("utf-8"))


def cli_muta(exe: Path) -> bool:
    """L'eseguibile parla ancora quando lo si lancia da un terminale?

    **E' la guardia che tiene insieme le due meta' del cambio.** Il pacchetto
    e' costruito senza console, cosi' il doppio click non apre nessuna finestra
    nera; ma quella scelta, da sola, rende `sys.stdout` `None` e trasforma
    `MrRao.exe health` in un comando che funziona **senza dire niente**. Non
    solleva, non fallisce, non lascia traccia: e' il modo peggiore di
    rompersi, e nessuno degli altri controlli qui dentro se ne accorgerebbe.

    Il controllo e' l'unico che possa dirlo: si lancia con un argomento, si
    legge cio' che arriva, e se non arriva niente il pacchetto viene respinto.

    Perche' `health` e non `--help`: perche' esercita anche l'aggancio *prima*
    che il resto del programma parta, che e' il punto fragile.
    """
    print("controllo che la riga di comando non sia diventata muta...")
    try:
        esito = subprocess.run(
            [str(exe), "health"],
            cwd=str(exe.parent),
            capture_output=True,
            timeout=120,
        )
    except subprocess.TimeoutExpired:
        print("FALLITO  'MrRao.exe health' non e' tornato entro 120 s",
              file=sys.stderr)
        return True

    uscita = (esito.stdout or b"") + (esito.stderr or b"")
    testo = uscita.decode("utf-8", "replace").strip()
    if not testo:
        print(
            "FALLITO  'MrRao.exe health' non ha stampato NIENTE.\n"
            "         Il pacchetto e' costruito senza console (--noconsole in\n"
            "         build_portable.bat) e l'aggancio in console_win.py non\n"
            "         ha funzionato: la riga di comando e' muta.",
            file=sys.stderr,
        )
        return True
    print(f"  OK     la CLI parla: {testo.splitlines()[0][:70]!r}")
    return False


def finestra_assente(exe: Path) -> bool:
    """Il pacchetto sa aprire la finestra dell'applicazione?

    **Il ripiego sul browser e' progettato per essere silenzioso**, ed e'
    giusto che lo sia per l'utente: su una macchina senza il motore di
    rendering di sistema il programma deve funzionare lo stesso. Ma la stessa
    qualita' lo rende invisibile in un pacchetto **mal costruito**: pywebview
    sceglie il backend a runtime, quindi se i suoi moduli non sono stati
    inclusi il pacchetto esce, funziona, apre il browser, e sembra una scelta.

    Qui si chiede all'eseguibile di dirlo, invece di dedurlo dalla dimensione
    del pacchetto o dalla presenza di un file.
    """
    print("controllo che la finestra dell'applicazione sia dentro il pacchetto...")
    try:
        esito = subprocess.run(
            [str(exe), "health"],
            cwd=str(exe.parent),
            capture_output=True,
            timeout=120,
        )
    except subprocess.TimeoutExpired:
        print("FALLITO  'MrRao.exe health' non e' tornato entro 120 s", file=sys.stderr)
        return True

    testo = ((esito.stdout or b"") + (esito.stderr or b"")).decode("utf-8", "replace")
    if "finestra: ok" in testo:
        print("  OK     la finestra e' disponibile")
        return False
    if "finestra: assente" in testo:
        print(
            "FALLITO  il pacchetto non sa aprire la finestra dell'applicazione.\n"
            "         Mancano i --hidden-import di webview in build_portable.bat:\n"
            "         il pacchetto ripiegherebbe sul browser in silenzio.",
            file=sys.stderr,
        )
        return True
    print(
        "FALLITO  'health' non dice niente sulla finestra: il controllo non sta\n"
        "         guardando piu' quello che credeva.",
        file=sys.stderr,
    )
    return True


def main(argv: list[str]) -> int:
    exe = Path(argv[1] if len(argv) > 1 else "dist/MrRao-Portable/app/MrRao.exe").resolve()
    if not exe.is_file():
        print(f"ERRORE: eseguibile assente: {exe}", file=sys.stderr)
        return 1

    porta = porta_libera()
    base = f"http://127.0.0.1:{porta}"
    ambiente = {
        **os.environ,
        "MR_RAO_PORT": str(porta),
        "MR_RAO_TRAY": "0",
        "MR_RAO_OPEN_BROWSER": "0",
        # **E la finestra spenta.** Senza, ogni verifica apriva una finestra
        # sullo schermo di chi costruisce, e — peggio — ne lasciava dietro il
        # processo: al build successivo la cartella `dist` risultava occupata,
        # la copia falliva a meta', e il pacchetto veniva respinto per un
        # motivo che non c'entrava niente con il pacchetto. Qui serve il
        # server, non l'interfaccia.
        "MR_RAO_FINESTRA": "0",
    }
    print(f"avvio  {exe.name} sulla porta {porta}")
    proc = subprocess.Popen(
        [str(exe)],
        cwd=str(exe.parent),
        env=ambiente,
        stdout=subprocess.PIPE,
        stderr=subprocess.STDOUT,
    )

    try:
        salute = None
        scaduto = time.monotonic() + AVVIO_MAX_S
        while time.monotonic() < scaduto:
            if proc.poll() is not None:
                uscita = proc.stdout.read().decode("utf-8", "replace") if proc.stdout else ""
                print(
                    f"FALLITO  il processo e' morto (codice {proc.returncode}) "
                    f"prima di rispondere.\n--- output ---\n{uscita[-3000:]}",
                    file=sys.stderr,
                )
                return 1
            try:
                with urllib.request.urlopen(f"{base}/api/health", timeout=3) as r:
                    salute = json.loads(r.read().decode("utf-8"))
                break
            except (urllib.error.URLError, OSError, TimeoutError):
                time.sleep(PASSO_S)

        if salute is None:
            print(f"FALLITO  nessuna risposta da /api/health in {AVVIO_MAX_S}s", file=sys.stderr)
            return 1

        print(f"  OK     health: v{salute.get('version')} frozen={salute.get('frozen')}")
        if not salute.get("frozen"):
            print("FALLITO  l'eseguibile non si dichiara frozen: non e' il pacchetto", file=sys.stderr)
            return 1

        # Conversioni vere, una per formato Office. Ognuno ha dietro una
        # libreria diversa: se ne manca una, il documento esce vuoto e
        # l'errore da' la colpa al file. Meglio che la colpa se la prenda
        # il build.
        campi = {
            "profile": "default",
            "privacy_filter": "true",
            "include_frontmatter": "false",
        }
        frase = "Contatta mario.rossi@example.it al 335 123 4567 in via Roma 12"
        prove = (
            # (estensione, contenuto, libreria dietro, etichette attese)
            (".docx", docx_di_prova(frase), "mammoth",
             ("EMAIL", "PHONE", "ADDRESS")),
            (".xlsx", xlsx_di_prova(frase), "pandas + openpyxl",
             ("EMAIL", "PHONE")),
            (".pptx", pptx_di_prova(frase), "python-pptx",
             ("EMAIL", "PHONE")),
        )
        for ext, contenuto, libreria, attesi in prove:
            esito = post_multipart(
                f"{base}/api/convert/sync", campi, f"verifica{ext}", contenuto
            )
            md = esito.get("markdown") or ""
            if "Contatta" not in md:
                print(
                    f"FALLITO  il {ext} non ha prodotto testo: nel pacchetto manca "
                    f"probabilmente {libreria}.\n{md[:400]}",
                    file=sys.stderr,
                )
                return 1
            mancanti = [e for e in attesi if not segnaposto_presente(md, e)]
            if mancanti:
                print(
                    f"FALLITO  anonimizzazione incompleta su {ext}, mancano "
                    f"{mancanti}:\n{md[:400]}",
                    file=sys.stderr,
                )
                return 1
            totale = (esito.get("redaction") or {}).get("total", 0)
            print(
                f"  OK     {ext:5s} {len(md):4d} caratteri, motore "
                f"{esito.get('engine')}, {totale} sostituzioni"
            )

        # L'icona spedita dev'essere quella del repository. Per un po' il
        # pacchetto ne ha portata una piu' povera, generata al passo 1 del
        # build da un percorso che sovrascriveva l'artwork rifinito a mano:
        # il collegamento sul Desktop funzionava, e mostrava l'icona
        # sbagliata. Nessuno se ne accorge guardando se il file c'e'.
        ico_repo = Path(__file__).resolve().parent.parent / "static" / "img" / "mr-rao.ico"
        # Portable Windows: <radice>/mr-rao.ico (exe in app/).
        # .app macOS: Contents/Resources/mr-rao.ico (exe in Contents/MacOS/).
        candidati_ico = (
            exe.parent.parent / "mr-rao.ico",
            exe.parent.parent / "Resources" / "mr-rao.ico",
            exe.parent / "mr-rao.ico",
        )
        ico_pacchetto = next((p for p in candidati_ico if p.is_file()), candidati_ico[0])
        if ico_pacchetto.is_file() and ico_repo.is_file():
            if ico_pacchetto.read_bytes() != ico_repo.read_bytes():
                print(
                    f"FALLITO  l'icona del pacchetto ({ico_pacchetto.stat().st_size:,} B) "
                    f"non e' quella del repository ({ico_repo.stat().st_size:,} B)",
                    file=sys.stderr,
                )
                return 1
            print(f"  OK     icona identica al repository: {ico_repo.stat().st_size:,} B")
        else:
            print(f"FALLITO  icona assente nel pacchetto: {ico_pacchetto}", file=sys.stderr)
            return 1

        if cli_muta(exe):
            return 1

        if finestra_assente(exe):
            return 1

        print("VERIFICA SUPERATA")
        return 0
    finally:
        proc.terminate()
        try:
            proc.wait(timeout=15)
        except subprocess.TimeoutExpired:
            proc.kill()


if __name__ == "__main__":
    raise SystemExit(main(sys.argv))
